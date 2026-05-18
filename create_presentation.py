"""Generate PowerPoint presentation for DarkStore Analytics Platform."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors as hex tuples (R, G, B)
DARK_BLUE = (0, 51, 102)
LIGHT_BLUE = (0, 112, 192)
ORANGE = (255, 102, 0)
WHITE = (255, 255, 255)
GRAY = (89, 89, 89)

def set_rgb(color_obj, rgb_tuple):
    """Set RGB color from tuple."""
    from pptx.dml.color import RGBColor
    color_obj.rgb = RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])


def add_title_slide(title, subtitle=""):
    from pptx.dml.color import RGBColor
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(DARK_BLUE[0], DARK_BLUE[1], DARK_BLUE[2])
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(GRAY[0], GRAY[1], GRAY[2])
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(title, bullets, subbullets=None):
    from pptx.dml.color import RGBColor
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(DARK_BLUE[0], DARK_BLUE[1], DARK_BLUE[2])

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(GRAY[0], GRAY[1], GRAY[2])
        p.space_before = Pt(12)
        p.level = 0

        # Add subbullets if provided
        if subbullets and i < len(subbullets) and subbullets[i]:
            for sub in subbullets[i]:
                p = tf.add_paragraph()
                p.text = sub
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(GRAY[0], GRAY[1], GRAY[2])
                p.level = 1

    return slide


def add_table_slide(title, headers, rows):
    from pptx.dml.color import RGBColor
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(DARK_BLUE[0], DARK_BLUE[1], DARK_BLUE[2])

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table_width = Inches(12)
    table_height = Inches(0.4 * num_rows)
    left = Inches(0.667)
    top = Inches(1.4)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table

    # Set column widths
    col_width = Inches(12 / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(DARK_BLUE[0], DARK_BLUE[1], DARK_BLUE[2])
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(WHITE[0], WHITE[1], WHITE[2])
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(GRAY[0], GRAY[1], GRAY[2])
            p.alignment = PP_ALIGN.LEFT

    return slide


def add_diagram_slide(title, content_text):
    from pptx.dml.color import RGBColor
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(DARK_BLUE[0], DARK_BLUE[1], DARK_BLUE[2])

    # Diagram as text (monospace)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = content_text
    p.font.size = Pt(11)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(GRAY[0], GRAY[1], GRAY[2])

    return slide


# ============== CREATE SLIDES ==============

# Slide 1: Title
add_title_slide(
    "DarkStore Analytics Platform",
    "An AI-Powered Microservices Architecture for Warehouse Operations"
)

# Slide 2: What is a Dark Store?
add_content_slide(
    "What is a Dark Store?",
    [
        "Definition: A retail distribution center that exclusively fulfills online orders",
        "Not open to public — optimized purely for operations",
        "Key Characteristics:",
        "Use Cases:"
    ],
    [
        None,
        None,
        ["Optimized for rapid order picking and packing", "Layout designed for operational efficiency", "Focus on speed-to-delivery metrics"],
        ["Grocery delivery (10-30 minute delivery)", "Quick commerce (q-commerce)", "Last-mile fulfillment centers"]
    ]
)

# Slide 3: The Business Problem
add_table_slide(
    "The Business Problem",
    ["Challenge", "Impact"],
    [
        ["Inefficient Picking Routes", "Workers walk unnecessary distances"],
        ["Reactive Inventory Management", "Stockouts discovered too late"],
        ["No Predictive Capabilities", "Cannot anticipate demand spikes"],
        ["Manual Anomaly Detection", "Issues go unnoticed"],
        ["Siloed Information", "Data trapped in separate systems"],
        ["Limited Scalability", "Systems cannot grow with business"],
    ]
)

# Slide 4: Solution Overview
add_content_slide(
    "Solution Overview: DarkStore Analytics Platform",
    [
        "A comprehensive analytics solution that provides:",
        "Real-time operational monitoring",
        "AI-powered demand forecasting",
        "Intelligent route optimization",
        "Automated anomaly detection",
        "Product affinity analysis",
        "Scenario simulation capabilities",
        "Natural language AI interface"
    ]
)

# Slide 5: Architecture Philosophy
add_content_slide(
    "Architecture Philosophy: Microservices",
    [
        "Why Microservices?",
        "Scalability — Scale individual services based on demand",
        "Resilience — Failure in one service doesn't crash the system",
        "Flexibility — Use different technologies per service",
        "Maintainability — Smaller codebases are easier to manage",
        "Deployment — Deploy services independently"
    ]
)

# Slide 6: System Architecture
add_diagram_slide(
    "System Architecture",
    """
                    ┌─────────────────────────────────────────────────┐
                    │              Client Applications                 │
                    │    (Dashboards, Mobile Apps, AI Assistants)     │
                    └─────────────────────────────────────────────────┘
                                            │
                              ┌─────────────┼─────────────┐
                              ▼             ▼             ▼
                        ┌──────────┐  ┌──────────┐  ┌──────────┐
                        │   API    │  │ Chatbot  │  │   MCP    │
                        │ Gateway  │  │ Service  │  │  Server  │
                        └────┬─────┘  └──────────┘  └──────────┘
                             │
                             ▼
                    ┌─────────────────────────────────────────────────┐
                    │            Backend Microservices                 │
                    ├─────────┬─────────┬─────────┬─────────┬────────┤
                    │  Data   │Forecast │ Route   │ Anomaly │Affinity│
                    │ Service │ Service │Optimizer│ Detector│Analyzer│
                    └─────────┴─────────┴─────────┴─────────┴────────┘
                                            │
                                            ▼
                                  ┌─────────────────┐
                                  │    Database     │
                                  └─────────────────┘
    """
)

# Slide 7: Component Overview
add_table_slide(
    "Component Overview: 9 Core Components",
    ["#", "Component", "Purpose"],
    [
        ["1", "API Gateway", "Central request routing"],
        ["2", "Data Service", "Core data access layer"],
        ["3", "Forecaster Service", "Demand prediction"],
        ["4", "Route Optimizer Service", "Picking path optimization"],
        ["5", "Anomaly Detector Service", "Pattern anomaly detection"],
        ["6", "Affinity Analyzer Service", "Product relationship analysis"],
        ["7", "Event Simulator Service", "Scenario modeling"],
        ["8", "AI Chatbot", "Natural language interface"],
        ["9", "MCP Server", "AI assistant integration"],
    ]
)

# Slide 8: API Gateway
add_content_slide(
    "API Gateway: The Central Entry Point",
    [
        "Purpose: Routes all incoming requests to the appropriate microservice",
        "Responsibilities:",
        "Request routing and load distribution",
        "Cross-Origin Resource Sharing (CORS) handling",
        "Health check aggregation across all services",
        "Unified API endpoint for clients",
        "Flow: Client Request → API Gateway → Microservice → Response"
    ]
)

# Slide 9: Data Service
add_content_slide(
    "Data Service: The Data Access Layer",
    [
        "Purpose: Provides unified access to all operational data",
        "Data Entities:",
        "Orders — Order details, status, fulfillment times",
        "Inventory — Stock levels, locations, expiry dates",
        "Alerts — System warnings and notifications",
        "Staff — Employee information and assignments",
        "Key Outputs: Summary KPIs, Low stock items, Active alerts, Staff metrics"
    ]
)

# Slide 10: Data Model
add_content_slide(
    "Data Model: Core Database Entities",
    [
        "Orders Table:",
        "Order ID, timestamp, items list, total value, quantity",
        "Status (Pending → Picking → Packing → Delivered)",
        "Assigned staff, fulfillment time, delay flag",
        "Inventory Table:",
        "SKU ID, product name, category, current stock, reorder threshold",
        "Shelf location (e.g., A05-3), unit price, expiry date",
        "Alerts Table:",
        "Alert type, severity, related SKU/order, status"
    ]
)

# Slide 11: KPIs
add_table_slide(
    "Key Performance Indicators (KPIs)",
    ["KPI", "Description"],
    [
        ["Total Orders", "Orders received (today / all-time)"],
        ["On-Time Rate", "Percentage meeting target fulfillment time"],
        ["Backlog", "Pending orders awaiting fulfillment"],
        ["Avg Fulfillment Time", "Mean time from order to dispatch"],
        ["Revenue", "Total and daily revenue figures"],
        ["Low Stock Count", "Items below reorder threshold"],
        ["Stock Health Rate", "Percentage of healthy inventory"],
        ["Active/Critical Alerts", "Open issues requiring attention"],
    ]
)

# Slide 12: Forecaster Service
add_content_slide(
    "Forecaster Service: Demand Prediction Engine",
    [
        "Purpose: Predicts future demand for each SKU",
        "Enables proactive inventory management",
        "Why Forecasting Matters:",
        "Prevent stockouts before they happen",
        "Optimize reorder quantities",
        "Plan for seasonal variations",
        "Reduce holding costs"
    ]
)

# Slide 13: Forecasting Models
add_table_slide(
    "Forecasting Models: Multi-Model Approach",
    ["Model", "Type", "Best For"],
    [
        ["ARIMA", "Statistical", "Stable demand patterns"],
        ["SARIMA", "Statistical", "Seasonal patterns"],
        ["Holt-Winters", "Statistical", "Trend + seasonality"],
        ["Random Forest", "Machine Learning", "Complex patterns"],
        ["Gradient Boosting", "Machine Learning", "High accuracy needs"],
        ["Ensemble", "Combined", "Balanced predictions"],
        ["Auto", "Automatic", "Selects best model"],
    ]
)

# Slide 14: Forecasting Process
add_content_slide(
    "Forecasting Process",
    [
        "Step 1: Collect Historical Data",
        "Step 2: Analyze Demand Patterns",
        "Step 3: Train Multiple Models",
        "Step 4: Generate Predictions",
        "Step 5: Calculate Reorder Recommendations",
        "Step 6: Output Forecast with Confidence Intervals",
        "Output: Daily predictions, reorder quantity, estimated cost, confidence metrics"
    ]
)

# Slide 15: Route Optimizer
add_content_slide(
    "Route Optimizer Service: Picking Path Optimization",
    [
        "Purpose: Calculates the most efficient route for warehouse pickers",
        "The Problem:",
        "Warehouses have thousands of shelf locations",
        "Pickers receive lists of items to collect",
        "Unoptimized routes waste time and energy",
        "Manual planning is impractical at scale",
        "Solution: Algorithmic route optimization"
    ]
)

# Slide 16: Warehouse Layout
add_content_slide(
    "Warehouse Layout Model",
    [
        "Location Format: A05-3",
        "A = Aisle letter (A-Z)",
        "05 = Bay number (01-20)",
        "3 = Shelf level (1-5)",
        "Coordinate System:",
        "Each location mapped to X,Y coordinates",
        "Distance calculated using Manhattan distance",
        "Special locations: ENTRANCE, PACKING"
    ]
)

# Slide 17: Route Algorithms
add_content_slide(
    "Route Optimization Algorithms",
    [
        "Two-Stage Optimization:",
        "Stage 1: Nearest Neighbor Algorithm",
        "Start at entrance, always move to closest unvisited location",
        "Greedy approach for quick initial solution",
        "Stage 2: 2-opt Improvement",
        "Take initial route, try swapping pairs of edges",
        "Keep improvements, discard worse solutions",
        "Repeat until no improvement found",
        "Result: Routes typically 30-45% shorter than naive ordering"
    ]
)

# Slide 18: Route Visualization
add_diagram_slide(
    "Route Optimization: Before vs After",
    """
        Naive Route:                         Optimized Route:

        ENTRANCE                             ENTRANCE
            ↓                                    ↓
           E15  (far)                           A05  (near)
            ↓                                    ↓
           A05  (backtrack)                     B12  (logical)
            ↓                                    ↓
           D08  (far again)                     C03  (sequential)
            ↓                                    ↓
           B12                                  D08
            ↓                                    ↓
           C03                                  E15
            ↓                                    ↓
        PACKING                              PACKING

        Distance: 330m                       Distance: 185m
                                             Improvement: 44%
    """
)

# Slide 19: Anomaly Detector
add_content_slide(
    "Anomaly Detector Service",
    [
        "Purpose: Identifies unusual patterns that may indicate problems",
        "Three Detection Areas:",
        "Order Anomalies — unusual order patterns",
        "Inventory Anomalies — stock discrepancies",
        "Staff Performance Anomalies — outlier behaviors"
    ]
)

# Slide 20: Anomaly Methods
add_content_slide(
    "Anomaly Detection Methods",
    [
        "Isolation Forest (for Orders):",
        "Unsupervised learning algorithm",
        "Isolates anomalies by random partitioning",
        "Effective for high-dimensional data",
        "Z-Score Analysis (for Inventory & Staff):",
        "Statistical method measuring standard deviations",
        "Flags values beyond threshold",
        "Simple and interpretable"
    ]
)

# Slide 21: Types of Anomalies
add_table_slide(
    "Types of Anomalies Detected",
    ["Category", "Anomaly Type", "Example"],
    [
        ["Orders", "Unusual value", "Order 10x normal size"],
        ["Orders", "Strange timing", "Orders at unusual hours"],
        ["Orders", "Fulfillment issues", "Extremely long processing"],
        ["Inventory", "Stock discrepancy", "Negative stock values"],
        ["Inventory", "Unusual movement", "Sudden stock changes"],
        ["Staff", "Performance outlier", "Unusually slow fulfillment"],
    ]
)

# Slide 22: Affinity Analyzer
add_content_slide(
    "Affinity Analyzer Service",
    [
        "Purpose: Discovers products frequently purchased together",
        "Business Applications:",
        "Optimize warehouse layout (place related items nearby)",
        "Bundle recommendations",
        "Cross-selling opportunities",
        "Demand correlation understanding"
    ]
)

# Slide 23: Affinity Metrics
add_table_slide(
    "Affinity Metrics",
    ["Metric", "Definition", "Use"],
    [
        ["Support", "How often items appear together", "Frequency measure"],
        ["Confidence", "P(B|A) - If A bought, likelihood of B", "Prediction strength"],
        ["Lift", "How much more likely than random", "Relationship strength"],
    ]
)

# Slide 24: Co-location
add_content_slide(
    "Co-location Recommendations",
    [
        "Concept: Place frequently co-purchased products near each other",
        "Benefits:",
        "Reduces picker travel time",
        "Speeds up order fulfillment",
        "Improves operational efficiency",
        "Output:",
        "Product pairs with high affinity scores",
        "Current vs recommended shelf locations"
    ]
)

# Slide 25: Event Simulator
add_content_slide(
    "Event Simulator Service",
    [
        "Purpose: Simulates different operational scenarios",
        "Available Scenarios:",
        "Normal — Baseline operations",
        "Flash Sale — 3x order volume spike",
        "Supply Delay — Extended fulfillment times",
        "Use Cases: Capacity planning, Staff scheduling, Stress testing"
    ]
)

# Slide 26: AI Chatbot
add_content_slide(
    "AI Chatbot: Natural Language Interface",
    [
        "Purpose: Query the system using everyday language",
        "Powered By:",
        "LangChain — AI application framework",
        "LangGraph — Agent orchestration",
        "Google Gemini 2.0 Flash — Large Language Model",
        "Key Feature: Conversation memory maintains context"
    ]
)

# Slide 27: Chatbot Architecture
add_diagram_slide(
    "Chatbot Architecture",
    """
                    User Query (Natural Language)
                               ↓
                        LangChain Agent
                               ↓
                        Tool Selection
                               ↓
                    API Call to Services
                               ↓
                    Response Processing
                               ↓
                    Natural Language Answer
    """
)

# Slide 28: Chatbot Tools
add_table_slide(
    "Chatbot Tools: 8 Available Functions",
    ["Tool", "Function"],
    [
        ["get_kpis", "Retrieve summary KPIs"],
        ["get_low_stock", "Get items needing restock"],
        ["forecast_demand", "Predict SKU demand"],
        ["optimize_route", "Calculate optimal picking path"],
        ["detect_anomalies", "Find unusual patterns"],
        ["get_product_pairs", "Get co-purchased items"],
        ["get_alerts", "Retrieve active alerts"],
        ["get_staff_performance", "Get staff metrics"],
    ]
)

# Slide 29: Chatbot Examples
add_table_slide(
    "Example Chatbot Interactions",
    ["User Says", "System Does"],
    [
        ["What's our on-time rate?", "Calls get_kpis, extracts rate"],
        ["Show me low stock items", "Calls get_low_stock, formats list"],
        ["Forecast demand for SKU01000", "Calls forecast_demand with SKU"],
        ["Optimize route for A05, B12, C03", "Calls optimize_route, shows path"],
        ["Are there any anomalies?", "Calls detect_anomalies, summarizes"],
        ["Which products are bought together?", "Calls get_product_pairs"],
    ]
)

# Slide 30: MCP Server
add_content_slide(
    "MCP Server: AI Assistant Integration",
    [
        "What is MCP?",
        "Model Context Protocol — standard for AI to interact with tools",
        "Purpose: Enables Claude AI to directly analyze dark store data",
        "12 Tools Available:",
        "Forecasting, route optimization, anomaly detection",
        "Product affinity, simulation, KPIs, and more"
    ]
)

# Slide 31: MCP Tools
add_table_slide(
    "MCP Server: 12 Tools for AI Integration",
    ["#", "Tool", "Purpose"],
    [
        ["1", "forecast_demand", "SKU demand prediction"],
        ["2", "compare_forecast_models", "Model comparison"],
        ["3", "optimize_picking_route", "Route optimization"],
        ["4", "batch_optimize_routes", "Multiple routes"],
        ["5", "detect_order_anomalies", "Order analysis"],
        ["6", "detect_inventory_anomalies", "Stock anomalies"],
        ["7", "detect_staff_anomalies", "Staff analysis"],
        ["8", "get_product_affinities", "Co-purchase analysis"],
        ["9", "get_colocation_recommendations", "Layout suggestions"],
        ["10", "run_simulation", "Scenario simulation"],
        ["11", "get_dashboard_kpis", "KPI retrieval"],
        ["12", "get_low_stock_items", "Low stock identification"],
    ]
)

# Slide 32: Technology Stack
add_table_slide(
    "Technology Stack",
    ["Layer", "Technology"],
    [
        ["Backend Framework", "FastAPI + Uvicorn"],
        ["Database", "SQLite (dev) / PostgreSQL (prod)"],
        ["Data Processing", "Pandas, NumPy"],
        ["Machine Learning", "scikit-learn, statsmodels"],
        ["AI Chatbot", "LangChain, LangGraph"],
        ["LLM", "Google Gemini 2.0 Flash"],
        ["MCP Server", "FastMCP"],
        ["Data Validation", "Pydantic"],
        ["Containerization", "Docker"],
    ]
)

# Slide 33: Why FastAPI
add_content_slide(
    "Why FastAPI?",
    [
        "High Performance — One of the fastest Python frameworks",
        "Automatic Documentation — Swagger UI generated automatically",
        "Type Safety — Pydantic integration for validation",
        "Async Support — Native asynchronous request handling",
        "Easy to Learn — Clean, intuitive syntax"
    ]
)

# Slide 34: Why Pydantic
add_content_slide(
    "Why Pydantic?",
    [
        "Purpose: Ensures data integrity throughout the system",
        "Benefits:",
        "Automatic request/response validation",
        "Clear data schemas",
        "Type hints for IDE support",
        "Serialization to JSON",
        "Error messages for invalid data",
        "Example Models: Order, InventoryItem, Alert, ForecastResponse"
    ]
)

# Slide 35: Why LangChain
add_content_slide(
    "Why LangChain?",
    [
        "Purpose: Simplifies building LLM-powered applications",
        "Components Used:",
        "Tools — Define functions the AI can call",
        "Agents — Autonomous decision-making",
        "Memory — Conversation history",
        "LangGraph — Agent orchestration",
        "Benefit: Abstracts complexity of LLM integration"
    ]
)

# Slide 36: Data Flow KPI
add_diagram_slide(
    "Data Flow Example: KPI Request",
    """
            1. Client sends GET /api/v1/data/kpis
                               ↓
            2. API Gateway receives request
                               ↓
            3. Request forwarded to Data Service
                               ↓
            4. Data Service queries database
                               ↓
            5. Calculations performed (on-time rate, etc.)
                               ↓
            6. Response formatted as JSON
                               ↓
            7. Returned through API Gateway to client
    """
)

# Slide 37: Data Flow Forecast
add_diagram_slide(
    "Data Flow Example: Forecast Request",
    """
            1. Client sends POST /api/v1/forecast/sku
               Body: { sku_id: "SKU01000", forecast_days: 14 }
                               ↓
            2. API Gateway routes to Forecaster Service
                               ↓
            3. Forecaster retrieves historical order data
                               ↓
            4. Multiple models trained on data
                               ↓
            5. Best model selected (or ensemble used)
                               ↓
            6. Predictions generated for next 14 days
                               ↓
            7. Reorder recommendations calculated
                               ↓
            8. Response returned with forecast + recommendations
    """
)

# Slide 38: Deployment
add_diagram_slide(
    "Deployment Architecture",
    """
            ┌─────────────────────────────────────────────┐
            │              Docker Compose                  │
            ├─────────────────────────────────────────────┤
            │  ┌─────────┐ ┌─────────┐ ┌─────────┐       │
            │  │ Gateway │ │  Data   │ │Forecast │       │
            │  │ :8000   │ │ :8001   │ │ :8002   │       │
            │  └─────────┘ └─────────┘ └─────────┘       │
            │  ┌─────────┐ ┌─────────┐ ┌─────────┐       │
            │  │ Routes  │ │ Anomaly │ │Affinity │       │
            │  │ :8003   │ │ :8004   │ │ :8005   │       │
            │  └─────────┘ └─────────┘ └─────────┘       │
            │  ┌─────────┐ ┌─────────┐ ┌─────────┐       │
            │  │Simulate │ │ Chatbot │ │   MCP   │       │
            │  │ :8006   │ │ :8020   │ │ :8010   │       │
            │  └─────────┘ └─────────┘ └─────────┘       │
            ├─────────────────────────────────────────────┤
            │              Internal Network                │
            └─────────────────────────────────────────────┘
    """
)

# Slide 39: Project Structure
add_diagram_slide(
    "Project Structure",
    """
            darkstore-analytics-platform/
            │
            ├── services/
            │   ├── api-gateway/
            │   ├── data-service/
            │   ├── forecaster-service/
            │   ├── route-optimizer-service/
            │   ├── anomaly-detector-service/
            │   ├── affinity-analyzer-service/
            │   └── event-simulator-service/
            │
            ├── chatbot/
            ├── mcp-server/
            ├── shared/          (Common Pydantic models)
            ├── data/            (SQLite database)
            └── docker-compose.yml
    """
)

# Slide 40: Scalability
add_content_slide(
    "Scalability",
    [
        "Horizontal Scaling:",
        "Each microservice can be scaled independently",
        "Add more instances of busy services",
        "Load balancer distributes requests",
        "Vertical Scaling:",
        "Increase resources for specific services",
        "Forecaster may need more CPU for ML",
        "Database Scaling:",
        "SQLite → PostgreSQL for production",
        "Database replication for read scaling"
    ]
)

# Slide 41: Benefits
add_table_slide(
    "Benefits Summary",
    ["Benefit", "Description"],
    [
        ["Operational Visibility", "Real-time KPIs and alerts"],
        ["Predictive Capability", "Demand forecasting prevents stockouts"],
        ["Efficiency Gains", "Optimized routes save 30-45% travel"],
        ["Proactive Detection", "Anomalies caught automatically"],
        ["Better Layout", "Affinity analysis improves warehouse"],
        ["Accessible Insights", "Natural language queries for everyone"],
        ["AI Integration", "Claude AI can analyze operations"],
        ["Future Ready", "Microservices architecture scales"],
    ]
)

# Slide 42: Use Cases
add_table_slide(
    "Use Cases: Who Uses What",
    ["User Role", "Primary Features"],
    [
        ["Warehouse Manager", "KPIs, Alerts, Staff Performance"],
        ["Inventory Planner", "Forecasting, Low Stock, Reorder Recommendations"],
        ["Operations Lead", "Route Optimization, Anomaly Detection"],
        ["Layout Designer", "Affinity Analysis, Co-location Recommendations"],
        ["Capacity Planner", "Event Simulation"],
        ["Executive", "Chatbot for quick insights"],
    ]
)

# Slide 43: Summary
add_content_slide(
    "Summary",
    [
        "DarkStore Analytics Platform transforms warehouse operations through:",
        "Microservices Architecture — Scalable, maintainable, resilient",
        "AI-Powered Analytics — Forecasting, optimization, detection",
        "Natural Language Interface — Accessible to all users",
        "Real-time Insights — KPIs, alerts, and recommendations",
        "Modern Technology Stack — FastAPI, LangChain, Docker"
    ]
)

# Slide 44: Thank You
add_title_slide(
    "Thank You",
    "7 Microservices • AI Chatbot • MCP Server • 12 MCP Tools • 8 Chatbot Tools • 15+ API Endpoints"
)

# Save
prs.save('/Users/ash/Projects/darkstore-analytics-platform/DarkStore_Analytics_Presentation.pptx')
print("Presentation saved to DarkStore_Analytics_Presentation.pptx")
